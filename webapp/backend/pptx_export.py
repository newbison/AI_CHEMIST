"""Markdown → PowerPoint (.pptx) 转换模块。

两条路径：
1. PptxGenJS 路径（主力）：DeepSeek 生成 JS 代码 → Node.js 执行 → 高质量 PPT
2. python-pptx 路径（回退）：解析 Markdown 直接用 python-pptx 渲染

PptxGenJS 生成的 PPT 设计质量更高，支持丰富的形状、阴影、图标等效果。
"""
from __future__ import annotations

import io
import math
import os
import re
import subprocess
import tempfile
from datetime import datetime
from pathlib import Path
from typing import Literal

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# 颜色主题（与 docx_export 一致）
COLOR_PRIMARY = RGBColor(0x1A, 0x3C, 0x6E)      # 深蓝
COLOR_SECONDARY = RGBColor(0x2E, 0x5C, 0x8A)    # 中蓝
COLOR_ACCENT = RGBColor(0x0D, 0x94, 0x88)       # 青绿
COLOR_TEXT = RGBColor(0x1F, 0x29, 0x37)         # 正文
COLOR_MUTED = RGBColor(0x6B, 0x72, 0x80)        # 次要
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT_BG = RGBColor(0xF0, 0xF4, 0xF8)     # 浅蓝灰背景
COLOR_TABLE_HEADER = RGBColor(0x1A, 0x3C, 0x6E)
COLOR_TABLE_ZEBRA = RGBColor(0xF0, 0xF4, 0xF8)

import os
_CN_FONT = os.environ.get("CN_FONT", "Noto Sans CJK SC" if os.name != "nt" else "Microsoft YaHei")
FONT_HEADER = _CN_FONT
FONT_BODY = _CN_FONT


def _strip_md(text: str) -> str:
    """去除 Markdown 行内标记，返回纯文本。"""
    text = re.sub(r'\*\*([^*]+)\*\*', r'\1', text)
    text = re.sub(r'`([^`]+)`', r'\1', text)
    text = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', text)
    return text.strip()


def _parse_table_row(line: str) -> list[str]:
    line = line.strip()
    if line.startswith('|'):
        line = line[1:]
    if line.endswith('|'):
        line = line[:-1]
    return [cell.strip() for cell in line.split('|')]


def _is_table_separator(line: str) -> bool:
    return bool(re.match(r'^\s*\|?[\s\-:]+(\|[\s\-:]+)+\|?\s*$', line))


def _add_cover_slide(prs: Presentation, title: str, subtitle: str) -> None:
    """封面页：深蓝背景 + 居中标题。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # 背景
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRIMARY

    # 标题
    left = Inches(0.8)
    top = Inches(2.0)
    width = Inches(8.4)
    height = Inches(1.5)
    tx_box = slide.shapes.add_textbox(left, top, width, height)
    tf = tx_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.name = FONT_HEADER
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = COLOR_WHITE

    # 副标题
    top2 = Inches(3.7)
    tx_box2 = slide.shapes.add_textbox(left, top2, width, Inches(0.8))
    tf2 = tx_box2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = subtitle
    run2.font.name = FONT_HEADER
    run2.font.size = Pt(16)
    run2.font.color.rgb = RGBColor(0xCAD, 0xCFC, 0xFC) if False else RGBColor(0xCA, 0xDC, 0xFC)

    # 日期
    top3 = Inches(4.8)
    tx_box3 = slide.shapes.add_textbox(left, top3, width, Inches(0.5))
    tf3 = tx_box3.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = datetime.now().strftime("%Y-%m-%d")
    run3.font.name = FONT_BODY
    run3.font.size = Pt(12)
    run3.font.color.rgb = RGBColor(0xAA, 0xBB, 0xCC)

    # 装饰线
    line_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(3.5), Inches(3.4), Inches(3.0), Pt(3)
    )
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = COLOR_ACCENT
    line_shape.line.fill.background()


def _add_section_slide(prs: Presentation, title: str) -> None:
    """章节分隔页：浅背景 + 居中标题。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_LIGHT_BG

    # 左侧色条
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(0.15), Inches(5.625)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()

    # 标题
    tx_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.4), Inches(1.2))
    tf = tx_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.name = FONT_HEADER
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = COLOR_PRIMARY


def _add_toc_slide(prs: Presentation, toc_items: list[tuple[str, list[str]]]) -> None:
    """目录检索页：列出所有章节（H1）及子章节（H2），作为导航检索。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题栏
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(0.8)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = COLOR_PRIMARY
    title_bar.line.fill.background()

    tx_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.1), Inches(9.2), Inches(0.6))
    tf = tx_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "目录 · CONTENTS"
    run.font.name = FONT_HEADER
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = COLOR_WHITE

    # 装饰线
    deco = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.4), Inches(0.85), Inches(2.0), Pt(2)
    )
    deco.fill.solid()
    deco.fill.fore_color.rgb = COLOR_ACCENT
    deco.line.fill.background()

    # 目录内容
    content_top = Inches(1.1)
    content_left = Inches(0.6)
    content_width = Inches(8.8)
    content_height = Inches(4.3)

    tx_box2 = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
    tf2 = tx_box2.text_frame
    tf2.word_wrap = True

    first = True
    for idx, (section_title, sub_items) in enumerate(toc_items, 1):
        if first:
            p = tf2.paragraphs[0]
            first = False
        else:
            p = tf2.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = f"{idx:02d}    {section_title}"
        run.font.name = FONT_HEADER
        run.font.size = Pt(15)
        run.font.bold = True
        run.font.color.rgb = COLOR_PRIMARY

        # 子章节（最多 5 个）
        for sub in sub_items[:5]:
            p = tf2.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(1)
            run = p.add_run()
            run.text = f"            ·  {sub}"
            run.font.name = FONT_BODY
            run.font.size = Pt(10)
            run.font.color.rgb = COLOR_MUTED


def _render_auto_scaled_table(
    slide, table_data: dict, left, top, max_width_in: float, max_height_in: float
) -> float:
    """渲染表格并自动放缩以适应可用宽度和高度。

    - 根据列数和内容长度动态调整字体大小（最小 6pt）
    - 按内容比例分配列宽（最小 0.5 英寸）
    - 超出高度时截断行数并添加提示行
    返回表格实际高度（英寸）。
    """
    headers = table_data["headers"]
    rows = table_data["rows"]
    num_cols = len(headers)
    if num_cols == 0:
        return 0.0

    # 计算每列最大字符数
    col_max_chars: list[int] = []
    for j in range(num_cols):
        max_len = len(_strip_md(headers[j])) if j < len(headers) else 0
        for row in rows:
            if j < len(row):
                max_len = max(max_len, len(_strip_md(row[j])))
        col_max_chars.append(max(max_len, 1))

    total_chars = sum(col_max_chars)

    # 根据列数选择初始字体
    if num_cols <= 3:
        font_size = 11
    elif num_cols <= 5:
        font_size = 10
    elif num_cols <= 7:
        font_size = 9
    else:
        font_size = 8

    # 计算列宽（按内容比例，最小 0.5 英寸）
    min_col_w = 0.5
    raw_widths = [max(min_col_w, max_width_in * (c / total_chars)) for c in col_max_chars]
    total_raw = sum(raw_widths)
    col_widths = [w * max_width_in / total_raw for w in raw_widths]

    # 逐列估算：在当前字体大小下，每列内容会换行成几行
    # 每行可容纳字符数 ≈ 列宽(英寸) * 72 / font_size * 0.9（0.9 为内边距系数）
    # 目标：每列最多换行 3 行，否则缩小字体
    def _max_lines_for_col(col_idx: int, fs: int) -> int:
        chars_per_line = col_widths[col_idx] * (72.0 / fs) * 0.9
        if chars_per_line < 1:
            return 999
        return max(1, math.ceil(col_max_chars[col_idx] / chars_per_line))

    while font_size > 6:
        worst_lines = max(_max_lines_for_col(j, font_size) for j in range(num_cols))
        if worst_lines <= 3:
            break
        font_size -= 1

    # 重新计算行高：基于实际最大换行数
    max_lines = max(_max_lines_for_col(j, font_size) for j in range(num_cols))
    row_h = font_size * 1.4 * max_lines / 72.0
    header_h = font_size * 1.8 / 72.0

    # 计算可显示行数
    avail_for_rows = max_height_in - header_h
    max_display = max(3, int(avail_for_rows / row_h))

    truncated = len(rows) > max_display
    display_rows = rows[:max_display]

    total_table_rows = 1 + len(display_rows) + (1 if truncated else 0)
    table_h = header_h + len(display_rows) * row_h + (row_h if truncated else 0)

    table_shape = slide.shapes.add_table(
        total_table_rows, num_cols,
        left, top,
        Inches(max_width_in), Inches(table_h)
    )
    table = table_shape.table

    # 设置列宽
    for j, w in enumerate(col_widths):
        table.columns[j].width = Inches(w)

    # 表头
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = _strip_md(h)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_TABLE_HEADER
        cell.margin_top = Pt(2)
        cell.margin_bottom = Pt(2)
        cell.margin_left = Pt(3)
        cell.margin_right = Pt(3)
        tf = cell.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        para.font.size = Pt(font_size)
        para.font.bold = True
        para.font.color.rgb = COLOR_WHITE
        para.font.name = FONT_BODY

    # 数据行
    for ri, row in enumerate(display_rows):
        for j in range(num_cols):
            cell_text = row[j] if j < len(row) else ""
            cell = table.cell(ri + 1, j)
            cell.text = _strip_md(cell_text)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_TABLE_ZEBRA if ri % 2 == 1 else COLOR_WHITE
            cell.margin_top = Pt(1)
            cell.margin_bottom = Pt(1)
            cell.margin_left = Pt(3)
            cell.margin_right = Pt(3)
            tf = cell.text_frame
            tf.word_wrap = True
            para = tf.paragraphs[0]
            para.font.size = Pt(font_size)
            para.font.color.rgb = COLOR_TEXT
            para.font.name = FONT_BODY

    # 截断提示行
    if truncated:
        remaining = len(rows) - max_display
        for j in range(num_cols):
            cell = table.cell(total_table_rows - 1, j)
            if j == 0:
                cell.text = f"... 还有 {remaining} 行未显示"
            else:
                cell.text = ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_LIGHT_BG
            cell.margin_top = Pt(1)
            cell.margin_bottom = Pt(1)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(max(font_size - 1, 6))
            para.font.italic = True
            para.font.color.rgb = COLOR_MUTED
            para.font.name = FONT_BODY

    return table_h


def _add_content_slide(prs: Presentation, title: str, bullets: list[str], tables: list[dict]) -> None:
    """内容页：标题 + 要点列表 + 表格（表格自动放缩）。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题栏
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0), Inches(10), Inches(0.8)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = COLOR_PRIMARY
    title_bar.line.fill.background()

    tx_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.1), Inches(9.2), Inches(0.6))
    tf = tx_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.name = FONT_HEADER
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = COLOR_WHITE

    # 内容区（使用英寸数值便于计算）
    SLIDE_H = 5.625
    content_top = 1.1
    content_left = 0.5
    content_width = 9.0
    bottom_margin = 0.3

    current_top = content_top

    # 先放表格（如果有），再放要点
    if tables:
        avail_height = SLIDE_H - current_top - bottom_margin
        if bullets:
            # 有要点时表格最多占 55% 可用高度，留空间给要点
            avail_height = min(avail_height, (SLIDE_H - content_top - bottom_margin) * 0.55)

        for table_data in tables:
            table_h = _render_auto_scaled_table(
                slide, table_data,
                Inches(content_left), Inches(current_top),
                content_width, avail_height
            )
            current_top += table_h + 0.2
            avail_height -= (table_h + 0.2)
            if avail_height < 0.5:
                break

    # 要点列表
    if bullets:
        bullet_height = SLIDE_H - current_top - bottom_margin
        if bullet_height < 0.5:
            bullet_height = 0.5
        tx_box = slide.shapes.add_textbox(
            Inches(content_left), Inches(current_top),
            Inches(content_width), Inches(bullet_height)
        )
        tf = tx_box.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(6)

            # 检测缩进层级
            indent_match = re.match(r'^(\s+)', bullet)
            indent = len(indent_match.group(1)) if indent_match else 0
            bullet_text = bullet.strip()
            # 去掉列表标记
            bullet_text = re.sub(r'^[-*+]\s+', '', bullet_text)
            bullet_text = re.sub(r'^\d+\.\s+', '', bullet_text)

            p.level = min(indent // 2, 2)

            run = p.add_run()
            run.text = _strip_md(bullet_text)
            run.font.name = FONT_BODY
            run.font.size = Pt(13) if p.level == 0 else Pt(11)
            run.font.color.rgb = COLOR_TEXT if p.level == 0 else COLOR_MUTED
            if p.level == 0:
                run.font.bold = False


def _add_closing_slide(prs: Presentation) -> None:
    """结尾页。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRIMARY

    tx_box = slide.shapes.add_textbox(Inches(1), Inches(2.3), Inches(8), Inches(1))
    tf = tx_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Thank You"
    run.font.name = FONT_HEADER
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = COLOR_WHITE

    tx_box2 = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.5))
    tf2 = tx_box2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "R&D Portfolio Intelligence Report"
    run2.font.name = FONT_BODY
    run2.font.size = Pt(14)
    run2.font.color.rgb = RGBColor(0xAA, 0xBB, 0xCC)


def markdown_to_pptx(markdown: str, *, title: str | None = None) -> bytes:
    """把 Markdown 报告转成 .pptx 字节流。"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # 提取标题
    if not title:
        first_heading = re.search(r'^#\s+(.+)$', markdown, re.MULTILINE)
        title = first_heading.group(1).strip() if first_heading else "R&D 智能报告"

    # 封面
    _add_cover_slide(prs, title, "R&D Portfolio Intelligence Report")

    # 预扫描：收集章节结构用于目录检索页
    # 策略：若存在多个 H1（除标题外），以 H1 为目录条目、H2 为子项；
    #       若仅有一个 H1（标题），则以 H2 为目录条目、H3 为子项。
    toc_items: list[tuple[str, list[str]]] = []
    _scan_lines = markdown.split('\n')
    _in_code = False
    _all_h1s: list[str] = []
    _all_h2s: list[str] = []
    _all_h3s: list[str] = []
    for _line in _scan_lines:
        if _line.strip().startswith('```'):
            _in_code = not _in_code
            continue
        if _in_code:
            continue
        _m = re.match(r'^(#{1,3})\s+(.*)', _line)
        if _m:
            _lvl = len(_m.group(1))
            _txt = _m.group(2).strip()
            if _lvl == 1:
                _all_h1s.append(_txt)
            elif _lvl == 2:
                _all_h2s.append(_txt)
            elif _lvl == 3:
                _all_h3s.append(_txt)

    # 判断用哪一级作为目录主条目
    _h1_sections = _all_h1s[1:] if len(_all_h1s) > 1 else []  # 排除第一个 H1（标题）

    if _h1_sections:
        # 以 H1 为目录条目，收集各自下面的 H2
        _cur_h1: str | None = None
        _cur_h2s: list[str] = []
        _h1_idx = 0
        _in_code2 = False
        for _line in _scan_lines:
            if _line.strip().startswith('```'):
                _in_code2 = not _in_code2
                continue
            if _in_code2:
                continue
            _m = re.match(r'^(#{1,2})\s+(.*)', _line)
            if _m:
                _lvl = len(_m.group(1))
                _txt = _m.group(2).strip()
                if _lvl == 1:
                    _h1_idx += 1
                    if _h1_idx == 1:
                        continue  # 跳过标题
                    if _cur_h1 is not None:
                        toc_items.append((_cur_h1, _cur_h2s))
                    _cur_h1 = _txt
                    _cur_h2s = []
                elif _lvl == 2 and _cur_h1 is not None:
                    _cur_h2s.append(_txt)
        if _cur_h1 is not None:
            toc_items.append((_cur_h1, _cur_h2s))
    else:
        # 以 H2 为目录条目，收集各自下面的 H3
        _cur_h2: str | None = None
        _cur_h3s: list[str] = []
        _in_code3 = False
        for _line in _scan_lines:
            if _line.strip().startswith('```'):
                _in_code3 = not _in_code3
                continue
            if _in_code3:
                continue
            _m = re.match(r'^(#{2,3})\s+(.*)', _line)
            if _m:
                _lvl = len(_m.group(1))
                _txt = _m.group(2).strip()
                if _lvl == 2:
                    if _cur_h2 is not None:
                        toc_items.append((_cur_h2, _cur_h3s))
                    _cur_h2 = _txt
                    _cur_h3s = []
                elif _lvl == 3 and _cur_h2 is not None:
                    _cur_h3s.append(_txt)
        if _cur_h2 is not None:
            toc_items.append((_cur_h2, _cur_h3s))

    # 目录检索页（封面之后、正文之前）
    if toc_items:
        _add_toc_slide(prs, toc_items)

    # 解析正文，按 H1/H2 分页
    lines = markdown.split('\n')
    i = 0
    skip_first_h1 = True
    current_section_title = ""
    current_slide_title = ""
    current_bullets: list[str] = []
    current_tables: list[dict] = []
    in_code_block = False
    code_buffer: list[str] = []

    def flush_content_slide():
        """把当前收集的 bullets/tables 输出为一页。"""
        nonlocal current_bullets, current_tables, current_slide_title
        if current_slide_title and (current_bullets or current_tables):
            _add_content_slide(prs, current_slide_title, current_bullets, current_tables)
        current_bullets = []
        current_tables = []
        current_slide_title = ""

    while i < len(lines):
        line = lines[i]

        # 代码块
        if line.strip().startswith('```'):
            if in_code_block:
                # 代码块结束，作为要点加入
                if code_buffer:
                    code_text = '\n'.join(code_buffer)
                    current_bullets.append(f"```\n{code_text}\n```")
                    code_buffer = []
                in_code_block = False
            else:
                in_code_block = True
            i += 1
            continue
        if in_code_block:
            code_buffer.append(line)
            i += 1
            continue

        # 空行
        if not line.strip():
            i += 1
            continue

        # 标题
        heading_match = re.match(r'^(#{1,4})\s+(.*)', line)
        if heading_match:
            level = len(heading_match.group(1))
            text = heading_match.group(2).strip()

            # 跳过第一个 H1（已用作封面）
            if level == 1 and skip_first_h1:
                skip_first_h1 = False
                i += 1
                continue

            # H1 = 章节分隔页
            if level == 1:
                flush_content_slide()
                _add_section_slide(prs, text)
                current_section_title = text
                i += 1
                continue

            # H2 = 新内容页标题
            if level == 2:
                flush_content_slide()
                current_slide_title = text
                i += 1
                continue

            # H3/H4 = 内容页内的子标题，作为要点
            if level >= 3:
                if current_slide_title:
                    current_bullets.append(f"{'  ' * (level - 3)}**{text}**")
                i += 1
                continue

        # 分隔线
        if line.strip() in ('---', '***', '___'):
            i += 1
            continue

        # 引用
        if line.strip().startswith('>'):
            text = re.sub(r'^>\s*', '', line.strip())
            if current_slide_title:
                current_bullets.append(f"  {text}")
            i += 1
            continue

        # 表格
        if '|' in line and i + 1 < len(lines) and _is_table_separator(lines[i + 1]):
            table_lines = [line]
            i += 1  # 跳过分隔行
            i += 1
            while i < len(lines) and '|' in lines[i] and lines[i].strip():
                table_lines.append(lines[i])
                i += 1
            headers = _parse_table_row(table_lines[0])
            rows = [_parse_table_row(r) for r in table_lines[1:]]
            current_tables.append({"headers": headers, "rows": rows})
            continue

        # 列表项
        list_match = re.match(r'^(\s*)([-*+]|\d+\.)\s+(.*)', line)
        if list_match:
            indent = len(list_match.group(1))
            text = list_match.group(3)
            if current_slide_title:
                current_bullets.append(f"{'  ' * (indent // 2)}{text}")
            i += 1
            continue

        # 普通段落
        if current_slide_title:
            current_bullets.append(line.strip())
        i += 1

    # 输出最后一页
    flush_content_slide()

    # 结尾页
    _add_closing_slide(prs)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ===========================================================================
# 设计驱动渲染 — 按设计方案 JSON 渲染美化版 PPT
# ===========================================================================

# accent 配色变体映射
ACCENT_COLORS: dict[str, RGBColor] = {
    "blue": RGBColor(0x2E, 0x5C, 0x8A),    # 中蓝
    "teal": RGBColor(0x0D, 0x94, 0x88),    # 青绿
    "amber": RGBColor(0xD9, 0x77, 0x06),   # 琥珀
    "slate": RGBColor(0x47, 0x55, 0x69),   # 石板灰
    "indigo": RGBColor(0x4F, 0x46, 0xE5),  # 靛蓝
}

ACCENT_LIGHT_BG: dict[str, RGBColor] = {
    "blue": RGBColor(0xEA, 0xF0, 0xF7),
    "teal": RGBColor(0xE6, 0xF7, 0xF5),
    "amber": RGBColor(0xFE, 0xF3, 0xE2),
    "slate": RGBColor(0xF1, 0xF3, 0xF5),
    "indigo": RGBColor(0xEE, 0xEA, 0xFD),
}


def _get_accent(name: str) -> RGBColor:
    return ACCENT_COLORS.get(name, COLOR_ACCENT)


def _get_accent_bg(name: str) -> RGBColor:
    return ACCENT_LIGHT_BG.get(name, COLOR_LIGHT_BG)


def _add_designed_cover(prs: Presentation, slide_spec: dict) -> None:
    """设计版封面页。"""
    title = slide_spec.get("title", "R&D 智能报告")
    subtitle = slide_spec.get("subtitle", "R&D Portfolio Intelligence Report")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRIMARY

    # 装饰几何块（右下角）
    deco = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(7.5), Inches(4.2), Inches(2.5), Inches(1.425)
    )
    deco.fill.solid()
    deco.fill.fore_color.rgb = COLOR_ACCENT
    deco.line.fill.background()

    # 标题
    tx = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(1.5))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.name = FONT_HEADER
    run.font.size = Pt(34)
    run.font.bold = True
    run.font.color.rgb = COLOR_WHITE

    # 副标题
    tx2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.4), Inches(8.4), Inches(0.6))
    p2 = tx2.text_frame.paragraphs[0]
    run2 = p2.add_run()
    run2.text = subtitle
    run2.font.name = FONT_BODY
    run2.font.size = Pt(15)
    run2.font.color.rgb = RGBColor(0xCA, 0xDC, 0xFC)

    # 日期
    tx3 = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(4), Inches(0.4))
    p3 = tx3.text_frame.paragraphs[0]
    run3 = p3.add_run()
    run3.text = datetime.now().strftime("%Y-%m-%d")
    run3.font.name = FONT_BODY
    run3.font.size = Pt(12)
    run3.font.color.rgb = RGBColor(0xAA, 0xBB, 0xCC)


def _add_designed_toc(prs: Presentation, slide_spec: dict) -> None:
    """设计版目录页。"""
    items = slide_spec.get("items", [])
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题栏
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(0.4), Inches(0.1), Inches(9.2), Inches(0.6))
    p = tx.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "目录 · CONTENTS"
    run.font.name = FONT_HEADER
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = COLOR_WHITE

    # 装饰线
    deco = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.85), Inches(2.0), Pt(2))
    deco.fill.solid()
    deco.fill.fore_color.rgb = COLOR_ACCENT
    deco.line.fill.background()

    # 目录项 — 双列布局
    tx2 = slide.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(8.8), Inches(4.3))
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    half = (len(items) + 1) // 2
    for idx, item in enumerate(items):
        p = tf2.paragraphs[0] if idx == 0 else tf2.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = f"{idx+1:02d}  {item}"
        run.font.name = FONT_HEADER
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = COLOR_PRIMARY


def _add_designed_section(prs: Presentation, slide_spec: dict) -> None:
    """设计版章节分隔页 — 带 accent 配色。"""
    title = slide_spec.get("title", "")
    accent_name = slide_spec.get("accent", "blue")
    accent = _get_accent(accent_name)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_LIGHT_BG

    # 左侧 accent 色条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.2), Inches(5.625))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    # 大编号圆
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.8), Inches(1.8), Inches(1.2), Inches(1.2)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent
    circle.line.fill.background()
    tf_c = circle.text_frame
    p_c = tf_c.paragraphs[0]
    p_c.alignment = PP_ALIGN.CENTER
    run_c = p_c.add_run()
    run_c.text = "§"
    run_c.font.name = FONT_HEADER
    run_c.font.size = Pt(28)
    run_c.font.bold = True
    run_c.font.color.rgb = COLOR_WHITE

    # 标题
    tx = slide.shapes.add_textbox(Inches(2.3), Inches(2.0), Inches(7.0), Inches(1.2))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = FONT_HEADER
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = COLOR_PRIMARY


def _add_callout_box(slide, text: str, left: float, top: float, width: float, accent: RGBColor) -> float:
    """添加高亮提示框，返回高度（英寸）。"""
    height = 0.8
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = _lighten(accent)
    box.line.color.rgb = accent
    box.line.width = Pt(1.5)

    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = FONT_BODY
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = accent
    return height


def _lighten(color: RGBColor) -> RGBColor:
    """把颜色变浅作为背景色。RGBColor 是 str 子类，如 '1A3C6E'。"""
    hex_str = str(color)
    r, g, b = int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16)
    return RGBColor(
        min(255, int(r + (255 - r) * 0.85)),
        min(255, int(g + (255 - g) * 0.85)),
        min(255, int(b + (255 - b) * 0.85)),
    )


def _add_designed_content(prs: Presentation, slide_spec: dict) -> None:
    """设计版内容页 — 支持 bullets/table/split/quote/comparison/metrics 布局。"""
    title = slide_spec.get("title", "")
    accent_name = slide_spec.get("accent", "blue")
    accent = _get_accent(accent_name)
    layout = slide_spec.get("layout", "bullets")
    highlights = slide_spec.get("highlights", [])
    callout = slide_spec.get("callout", "")
    bullets = slide_spec.get("bullets", [])
    tables = slide_spec.get("tables", [])

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题栏（带 accent 左侧色条）
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.7))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()

    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.08), Inches(0.7))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent
    accent_bar.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(0.3), Inches(0.08), Inches(9.4), Inches(0.55))
    p = tx.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = FONT_HEADER
    run.font.size = Pt(19)
    run.font.bold = True
    run.font.color.rgb = COLOR_WHITE

    SLIDE_H = 5.625
    content_top = 1.0
    content_left = 0.5
    content_width = 9.0
    bottom_margin = 0.3

    # ---- quote 布局：大字居中 ----
    if layout == "quote" and callout:
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.0), Inches(1.8), Inches(8.0), Inches(2.0)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = _get_accent_bg(accent_name)
        box.line.color.rgb = accent
        box.line.width = Pt(2)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(24)
        tf.margin_right = Pt(24)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = callout
        run.font.name = FONT_HEADER
        run.font.size = Pt(24)
        run.font.bold = True
        run.font.color.rgb = accent
        return

    # ---- metrics 布局：大数字卡片 ----
    if layout == "metrics" and bullets:
        card_data = [(b.get("text", ""), b.get("level", 0)) for b in bullets[:4]]
        card_w = 2.0
        gap = 0.2
        total_w = len(card_data) * card_w + (len(card_data) - 1) * gap
        start_left = (10.0 - total_w) / 2
        for i, (text, _) in enumerate(card_data):
            cl = start_left + i * (card_w + gap)
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(cl), Inches(1.8), Inches(card_w), Inches(1.8)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = _get_accent_bg(accent_name)
            card.line.color.rgb = accent
            card.line.width = Pt(1.5)
            tf = card.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = text
            run.font.name = FONT_HEADER
            run.font.size = Pt(22)
            run.font.bold = True
            run.font.color.rgb = accent
        return

    # ---- comparison 布局：左右对比 ----
    if layout == "comparison" and len(bullets) >= 2:
        mid = len(bullets) // 2
        left_bullets = bullets[:mid]
        right_bullets = bullets[mid:]
        col_w = 4.3
        for col_idx, (col_bullets, col_left, col_label) in enumerate([
            (left_bullets, 0.4, "A"),
            (right_bullets, 5.3, "B"),
        ]):
            # 列标题
            lbl = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(col_left), Inches(1.0), Inches(col_w), Inches(0.4)
            )
            lbl.fill.solid()
            lbl.fill.fore_color.rgb = accent if col_idx == 0 else COLOR_MUTED
            lbl.line.fill.background()
            tf = lbl.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = col_label
            run.font.name = FONT_HEADER
            run.font.size = Pt(13)
            run.font.bold = True
            run.font.color.rgb = COLOR_WHITE
            # 列内容
            tx = slide.shapes.add_textbox(Inches(col_left), Inches(1.6), Inches(col_w), Inches(3.5))
            tf = tx.text_frame
            tf.word_wrap = True
            for bi, b in enumerate(col_bullets):
                p = tf.paragraphs[0] if bi == 0 else tf.add_paragraph()
                p.space_after = Pt(6)
                run = p.add_run()
                run.text = b.get("text", "")
                run.font.name = FONT_BODY
                run.font.size = Pt(12)
                run.font.color.rgb = COLOR_TEXT
        return

    # ---- split 布局：左要点 + 右表格 ----
    if layout == "split" and (bullets or tables):
        left_w = 4.3
        right_left = 5.0
        right_w = 4.5
        # 左侧要点
        if bullets:
            tx = slide.shapes.add_textbox(Inches(0.4), Inches(1.0), Inches(left_w), Inches(4.0))
            tf = tx.text_frame
            tf.word_wrap = True
            for bi, b in enumerate(bullets):
                p = tf.paragraphs[0] if bi == 0 else tf.add_paragraph()
                p.space_after = Pt(8)
                level = b.get("level", 0)
                text = b.get("text", "")
                run = p.add_run()
                prefix = "• " if level == 0 else "  ◦ "
                run.text = prefix + text
                run.font.name = FONT_BODY
                run.font.size = Pt(12) if level == 0 else Pt(11)
                run.font.color.rgb = COLOR_TEXT if level == 0 else COLOR_MUTED
        # 右侧表格
        if tables:
            _render_auto_scaled_table(
                slide, tables[0],
                Inches(right_left), Inches(1.0),
                right_w, SLIDE_H - 1.0 - bottom_margin
            )
        return

    # ---- table 布局：表格为主 ----
    if layout == "table" and tables:
        avail_h = SLIDE_H - content_top - bottom_margin
        if callout:
            avail_h -= 1.0
        current_top = content_top
        for td in tables:
            th = _render_auto_scaled_table(
                slide, td,
                Inches(content_left), Inches(current_top),
                content_width, avail_h
            )
            current_top += th + 0.2
            avail_h -= (th + 0.2)
        if callout:
            _add_callout_box(slide, callout, content_left, current_top + 0.1, content_width, accent)
        return

    # ---- bullets 布局（默认）----
    current_top = content_top
    if bullets:
        avail_h = SLIDE_H - current_top - bottom_margin
        if callout:
            avail_h -= 1.0
        tx = slide.shapes.add_textbox(Inches(content_left), Inches(current_top), Inches(content_width), Inches(avail_h))
        tf = tx.text_frame
        tf.word_wrap = True
        for bi, b in enumerate(bullets):
            p = tf.paragraphs[0] if bi == 0 else tf.add_paragraph()
            p.space_after = Pt(8)
            level = b.get("level", 0)
            text = b.get("text", "")
            run = p.add_run()
            prefix = "● " if level == 0 else "  ◦ "
            run.text = prefix + text
            run.font.name = FONT_BODY
            run.font.size = Pt(14) if level == 0 else Pt(12)
            run.font.bold = level == 0
            # 高亮关键词
            if any(h in text for h in highlights):
                run.font.color.rgb = accent
                run.font.bold = True
            else:
                run.font.color.rgb = COLOR_TEXT if level == 0 else COLOR_MUTED
        current_top += min(avail_h, len(bullets) * 0.4)

    if callout:
        _add_callout_box(slide, callout, content_left, current_top + 0.15, content_width, accent)


def _add_designed_closing(prs: Presentation) -> None:
    """设计版结尾页。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_PRIMARY

    # 装饰
    deco = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.2), Inches(10), Inches(0.425))
    deco.fill.solid()
    deco.fill.fore_color.rgb = COLOR_ACCENT
    deco.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(8), Inches(1))
    p = tx.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Thank You"
    run.font.name = FONT_HEADER
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = COLOR_WHITE

    tx2 = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(0.5))
    p2 = tx2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "R&D Portfolio Intelligence Report"
    run2.font.name = FONT_BODY
    run2.font.size = Pt(14)
    run2.font.color.rgb = RGBColor(0xAA, 0xBB, 0xCC)


def markdown_to_pptx_designed(design_plan: dict, *, title: str | None = None) -> bytes:
    """按设计方案 JSON 渲染美化版 PPT。

    Args:
        design_plan: ppt_designer.design_ppt_slides_safe() 返回的设计方案
        title: 可选覆盖标题

    Returns:
        .pptx 字节流。若设计方案为空，回退到普通 markdown_to_pptx。
    """
    slides = design_plan.get("slides", [])
    if not slides:
        # 回退：无设计方案时用普通渲染
        return b""

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    overall_title = title or design_plan.get("title", "R&D 智能报告")

    for spec in slides:
        slide_type = spec.get("type", "content")
        if slide_type == "cover":
            spec.setdefault("title", overall_title)
            _add_designed_cover(prs, spec)
        elif slide_type == "toc":
            _add_designed_toc(prs, spec)
        elif slide_type == "section":
            _add_designed_section(prs, spec)
        elif slide_type == "closing":
            _add_designed_closing(prs)
        else:  # content
            _add_designed_content(prs, spec)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ===========================================================================
# PptxGenJS 路径：DeepSeek 生成 JS 代码 → Node.js 执行 → .pptx
# ===========================================================================


def execute_pptxgenjs(js_code: str, *, timeout: int = 60) -> bytes | None:
    """执行 PptxGenJS JavaScript 代码，返回生成的 .pptx 字节流。

    Args:
        js_code: PptxGenJS JavaScript 代码（完整的 Node.js 脚本）
        timeout: Node.js 执行超时秒数

    Returns:
        .pptx 文件字节流，失败返回 None
    """
    # 检查 Node.js 是否可用
    try:
        result = subprocess.run(
            ["node", "--version"], capture_output=True, text=True, timeout=5
        )
        if result.returncode != 0:
            print("[pptx_export] Node.js 不可用")
            return None
        print(f"[pptx_export] Node.js {result.stdout.strip()}")
    except FileNotFoundError:
        print("[pptx_export] 未找到 Node.js，请安装 Node.js")
        return None
    except Exception as e:
        print(f"[pptx_export] 检查 Node.js 失败: {e}")
        return None

    # 构建环境变量：加入全局 node_modules 路径
    env = os.environ.copy()
    # 查找全局 npm 路径
    try:
        npm_root = subprocess.run(
            ["npm", "root", "-g"], capture_output=True, text=True, timeout=5
        )
        global_modules = npm_root.stdout.strip()
        if global_modules and Path(global_modules).is_dir():
            existing = env.get("NODE_PATH", "")
            env["NODE_PATH"] = f"{global_modules};{existing}" if existing else global_modules
    except Exception:
        pass

    # 在临时目录中执行（PptxGenJS 的 writeFile 会输出到当前工作目录）
    with tempfile.TemporaryDirectory() as tmpdir:
        tmpdir_path = Path(tmpdir)

        # 写入 JS 脚本
        script_path = tmpdir_path / "generate.js"
        script_path.write_text(js_code, encoding="utf-8")

        # 执行
        try:
            result = subprocess.run(
                ["node", str(script_path)],
                capture_output=True,
                text=True,
                timeout=timeout,
                cwd=str(tmpdir_path),
                env=env,
            )
        except subprocess.TimeoutExpired:
            print(f"[pptx_export] Node.js 执行超时 ({timeout}s)")
            return None
        except Exception as e:
            print(f"[pptx_export] Node.js 执行异常: {e}")
            return None

        # 检查输出
        stdout = result.stdout
        stderr = result.stderr

        if stderr:
            print(f"[pptx_export] Node.js stderr:\n{stderr[:2000]}")

        if "DONE:" not in stdout:
            print(f"[pptx_export] 脚本未正常完成，stdout:\n{stdout[:2000]}")
            # 仍然尝试读取 output.pptx（可能 writeFile 成功但 console.log 格式不同）

        output_path = tmpdir_path / "output.pptx"
        if output_path.exists():
            pptx_bytes = output_path.read_bytes()
            print(f"[pptx_export] 成功生成 .pptx: {len(pptx_bytes)} bytes ({len(pptx_bytes)//1024} KB)")
            return pptx_bytes

        print("[pptx_export] 未找到 output.pptx")
        return None


def markdown_to_pptx_via_pptxgenjs(markdown: str, *, title: str | None = None) -> bytes | None:
    """PptxGenJS 路径：Markdown → DeepSeek 生成 JS 代码 → Node.js 执行 → .pptx

    这是推荐的 PPT 生成路径，生成的 PPT 设计质量远高于 python-pptx 路径。

    Args:
        markdown: R&D 报告 Markdown 文本
        title: 可选覆盖标题

    Returns:
        .pptx 字节流，失败返回 None（调用方应回退到 markdown_to_pptx）
    """
    from ppt_designer import design_pptx_js_safe

    js_code = design_pptx_js_safe(markdown)
    if not js_code:
        print("[pptx_export] DeepSeek 未生成有效的 PptxGenJS 代码，回退")
        return None

    print(f"[pptx_export] DeepSeek 生成了 {len(js_code)} 字符的 PptxGenJS 代码")

    pptx_bytes = execute_pptxgenjs(js_code)
    return pptx_bytes


if __name__ == "__main__":
    # 测试
    sample = Path(__file__).parent.parent.parent / "docs" / "reports" / "2026-06-19-battery-tape-rd-report.md"
    if sample.exists():
        md = sample.read_text(encoding="utf-8")
        pptx_bytes = markdown_to_pptx(md)
        out = Path("_test_report.pptx")
        out.write_bytes(pptx_bytes)
        print(f"生成 {out.resolve()}, {len(pptx_bytes)} bytes, {len(pptx_bytes)//1024} KB")
    else:
        print(f"未找到测试报告: {sample}")
